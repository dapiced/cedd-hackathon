#!/usr/bin/env python3
"""
generate_slides.py — Generate the 404HarmNotFound hackathon presentation (12 slides)

Mila x Bell x Kids Help Phone Hackathon · March 23, 2026
Team: 404HarmNotFound

Usage:
    python generate_slides.py

Output:
    presentation_404HarmNotFound.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x44, 0x44, 0x44)
ROW_ALT = RGBColor(0xE8, 0xEE, 0xF6)  # light blue-gray for table rows

GREEN = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
LIGHT_BLUE = RGBColor(0xD6, 0xE9, 0xF8)

FONT_NAME = "Calibri"
OUTPUT_FILE = "presentation_404HarmNotFound.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Set text in a shape's text frame."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment


def add_paragraph(tf, text, font_size=16, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(4), space_after=Pt(4), level=0, font_name=FONT_NAME):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_GRAY, bold_items=None):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check if item has bold prefix (text before |)
        is_bold = bold_items and i in bold_items
        if isinstance(item, tuple):
            # (bold_part, normal_part)
            run_bold = p.add_run()
            run_bold.text = item[0]
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = color
            run_bold.font.bold = True
            run_bold.font.name = FONT_NAME
            run_normal = p.add_run()
            run_normal.text = item[1]
            run_normal.font.size = Pt(font_size)
            run_normal.font.color.rgb = color
            run_normal.font.bold = False
            run_normal.font.name = FONT_NAME
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = is_bold
            p.font.name = FONT_NAME

        p.space_before = Pt(6)
        p.space_after = Pt(4)
        # Bullet character
        if not isinstance(item, tuple):
            p.text = f"\u2022  {item}"
        else:
            # prepend bullet to the bold run
            run_bold.text = f"\u2022  {item[0]}"

    return txBox


def add_slide_number(slide, number):
    """Add slide number at bottom-right."""
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.45), Inches(0.6), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(11)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at the top of a content slide."""
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.15), fill_color=DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), SLIDE_WIDTH - Inches(1.4), Inches(0.65))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), SLIDE_WIDTH - Inches(1.4), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.font.name = FONT_NAME


def add_table(slide, left, top, width, rows_data, col_widths=None, header_color=DARK_BLUE,
              header_font_color=WHITE, font_size=14, row_height=Inches(0.4)):
    """Add a styled table to a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_NAME
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_font_color
                else:
                    paragraph.font.color.rgb = TEXT_GRAY
                paragraph.alignment = PP_ALIGN.LEFT

            # Cell fill
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

            # Vertical alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """Add a thin accent line."""
    return add_shape(slide, left, top, width, height, fill_color=color)


def add_quote_box(slide, left, top, width, height, text, font_size=16, color=DARK_BLUE):
    """Add an italicized quote in a light box."""
    box = add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{text}"'
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = FONT_NAME
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    return box


def add_alert_dots(slide, left, top):
    """Add the four colored alert level dots."""
    colors = [GREEN, YELLOW, ORANGE, RED]
    labels = ["Green", "Yellow", "Orange", "Red"]
    for i, (c, lbl) in enumerate(zip(colors, labels)):
        x = left + Inches(i * 1.6)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, top, Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        dot.shadow.inherit = False

        txBox = slide.shapes.add_textbox(x + Inches(0.42), top + Inches(0.02), Inches(1.0), Inches(0.35))
        p = txBox.text_frame.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = FONT_NAME
        p.font.bold = True


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_slide_01_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DARK_BLUE)

    # Top accent line
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6), fill_color=ACCENT_BLUE)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CEDD"
    p.font.size = Pt(60)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.65), Inches(10.3), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Conversational Emotional Drift Detection"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    p2.font.bold = False
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    # Accent line below subtitle
    add_shape(slide, Inches(4.5), Inches(3.55), Inches(4.3), Pt(3), fill_color=ACCENT_BLUE)

    # Hackathon info
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Mila Hackathon  \u00b7  AI Safety & Youth Mental Health"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x99, 0xAA, 0xBB)
    p3.font.name = FONT_NAME
    p3.alignment = PP_ALIGN.CENTER

    # Team name with shield-like styling
    team_box = add_rounded_rect(slide, Inches(4.2), Inches(4.7), Inches(4.9), Inches(0.65),
                                 fill_color=RGBColor(0x24, 0x3B, 0x5E), line_color=ACCENT_BLUE)
    tf_team = team_box.text_frame
    tf_team.word_wrap = True
    p_team = tf_team.paragraphs[0]
    p_team.alignment = PP_ALIGN.CENTER
    run = p_team.add_run()
    run.text = "Team  404HarmNotFound"
    run.font.size = Pt(22)
    run.font.color.rgb = GREEN
    run.font.bold = True
    run.font.name = FONT_NAME

    # Date
    txBox4 = slide.shapes.add_textbox(Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.4))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "March 23, 2026"
    p4.font.size = Pt(16)
    p4.font.color.rgb = MEDIUM_GRAY
    p4.font.name = FONT_NAME
    p4.alignment = PP_ALIGN.CENTER

    # Bottom accent
    add_shape(slide, Inches(0), SLIDE_HEIGHT - Pt(6), SLIDE_WIDTH, Pt(6), fill_color=GREEN)

    add_slide_number(slide, 1)


def build_slide_02_problem(prs):
    """Slide 2: The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "The Problem", "Why current chatbots fail youth in crisis")

    # Left column: bullet points
    items = [
        "Youth mental health chatbots see each message in isolation",
        "They miss the trajectory: shorter messages, fewer questions, loss of hope words",
        "A youth can spiral across 10 messages and the chatbot keeps responding the same way",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(6.0), Inches(2.5), items, font_size=16)

    # Quote
    add_quote_box(slide, Inches(0.7), Inches(3.9), Inches(6.0), Inches(0.65),
                  "The signal is in the trajectory, not in the words.", font_size=17)

    # Right column: stats in colored boxes
    stats = [
        ("46%", "of KHP youth identify as 2SLGBTQ+", ACCENT_BLUE),
        ("2x", "Suicide contacts among youth 13 and under\ndoubled in 4 years", RED),
        ("44%", "of 988 callers abandon before connecting", ORANGE),
        ("75%", "share something they've never told anyone else", GREEN),
    ]

    y = Inches(1.5)
    for number, desc, color in stats:
        box = add_rounded_rect(slide, Inches(7.4), y, Inches(5.2), Inches(1.15),
                                fill_color=LIGHT_GRAY, line_color=color)
        # Number
        num_box = slide.shapes.add_textbox(Inches(7.6), y + Inches(0.1), Inches(1.3), Inches(0.9))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(32)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(9.0), y + Inches(0.15), Inches(3.4), Inches(0.85))
        tf2 = desc_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_GRAY
        p2.font.name = FONT_NAME

        y += Inches(1.3)

    add_slide_number(slide, 2)


def build_slide_03_solution(prs):
    """Slide 3: Our Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Our Solution: CEDD", "Real-time safety layer beside any chatbot")

    # Key points
    items = [
        "Monitors emotional drift across the full conversation, not single messages",
        "4 alert levels: Green \u2192 Yellow \u2192 Orange \u2192 Red",
        "Bilingual native: French + English",
        "~0ms detection latency, $0 compute (ML only, LLM for response only)",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(2.3), items, font_size=15)

    # Alert dots
    add_alert_dots(slide, Inches(0.9), Inches(3.7))

    # Comparison table
    table_data = [
        ["Classic Approach", "CEDD"],
        ["Keyword detection", "67-feature trajectory analysis"],
        ["Single response mode", "4 adaptive response modes"],
        ["Stateless", "Cross-session tracking"],
        ["Black box", "Fully interpretable features"],
        ["English only", "Bilingual FR + EN"],
    ]
    add_table(slide, Inches(7.0), Inches(1.4), Inches(5.8), table_data,
              col_widths=[Inches(2.9), Inches(2.9)], font_size=14, row_height=Inches(0.45))

    add_slide_number(slide, 3)


def build_slide_04_architecture(prs):
    """Slide 4: Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Architecture", "Three-layer detection + response system")

    # Layer boxes (stacked vertically)
    layers = [
        ("Layer 1: Lexical Safety Rules", "Crisis keywords \u2192 Orange/Red floor (priority override)",
         RED, RGBColor(0xFD, 0xED, 0xED)),
        ("Layer 2: ML Model", "67 features \u2192 GradientBoosting \u2192 alert level prediction",
         ORANGE, RGBColor(0xFE, 0xF3, 0xE4)),
        ("Layer 3: Response Modulation", "Adaptive system prompt \u2192 LLM fallback chain \u2192 safe response",
         ACCENT_BLUE, LIGHT_BLUE),
    ]

    y = Inches(1.5)
    for title, desc, border_color, bg_color in layers:
        box = add_rounded_rect(slide, Inches(0.7), y, Inches(7.5), Inches(1.15),
                                fill_color=bg_color, line_color=border_color)
        # Title
        t_box = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.1), Inches(7.0), Inches(0.45))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BLUE
        p.font.bold = True
        p.font.name = FONT_NAME

        # Description
        d_box = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.52), Inches(7.0), Inches(0.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_GRAY
        p2.font.name = FONT_NAME

        # Arrow between layers
        if y < Inches(3.5):
            arrow_box = slide.shapes.add_textbox(Inches(4.2), y + Inches(1.15), Inches(0.5), Inches(0.35))
            pa = arrow_box.text_frame.paragraphs[0]
            pa.text = "\u25bc"
            pa.font.size = Pt(18)
            pa.font.color.rgb = MEDIUM_GRAY
            pa.alignment = PP_ALIGN.CENTER

        y += Inches(1.5)

    # Right side: LLM fallback chain
    fb_box = add_rounded_rect(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(3.8),
                               fill_color=LIGHT_GRAY, line_color=ACCENT_BLUE)

    fb_title = slide.shapes.add_textbox(Inches(9.0), Inches(1.6), Inches(3.6), Inches(0.4))
    p = fb_title.text_frame.paragraphs[0]
    p.text = "LLM Fallback Chain"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    chain_items = [
        "1. Cohere Command A",
        "2. Groq (Llama 3.3 70B)",
        "3. Gemini 2.5 Flash",
        "4. Claude Haiku",
        "5. Static emergency text",
    ]
    add_bullet_list(slide, Inches(9.1), Inches(2.15), Inches(3.5), Inches(2.5),
                    chain_items, font_size=14, color=TEXT_GRAY)

    # Bottom note
    note_box = slide.shapes.add_textbox(Inches(9.0), Inches(4.5), Inches(3.6), Inches(0.5))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Detection: ~0ms, $0\n(ML only, no LLM needed)"
    p.font.size = Pt(13)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 4)


def build_slide_05_features(prs):
    """Slide 5: Feature Engineering (67 features)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Feature Engineering", "67 features capturing emotional trajectory")

    # Left: 10 per-message features table
    feature_data = [
        ["#", "Feature", "Measures"],
        ["0", "word_count", "Message length (shorter = concerning)"],
        ["1", "punctuation_ratio", "Punctuation density"],
        ["2", "question_presence", "Contains a question? (0/1)"],
        ["3", "negative_score", "Negative word ratio (bilingual)"],
        ["4", "finality_score", "Finality words (\"end it all\")"],
        ["5", "hope_score", "Hope/future word ratio"],
        ["6", "length_delta", "Length change vs previous msg"],
        ["7", "negation_score", "Negated positives (\"can't cope\")"],
        ["8", "identity_conflict", "2SLGBTQ+/cultural distress"],
        ["9", "somatization", "Physical + emotional co-occurrence"],
    ]
    add_table(slide, Inches(0.5), Inches(1.35), Inches(7.2), feature_data,
              col_widths=[Inches(0.4), Inches(2.0), Inches(4.8)],
              font_size=12, row_height=Inches(0.38))

    # Right: computation breakdown
    calc_box = add_rounded_rect(slide, Inches(8.2), Inches(1.35), Inches(4.6), Inches(4.5),
                                 fill_color=LIGHT_GRAY, line_color=ACCENT_BLUE)

    calc_title = slide.shapes.add_textbox(Inches(8.4), Inches(1.5), Inches(4.2), Inches(0.4))
    p = calc_title.text_frame.paragraphs[0]
    p.text = "Feature Computation"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    breakdown = [
        ("10 features", " \u00d7 6 temporal stats"),
        ("", " (mean, std, slope, last, max, min)"),
        ("= 60", " trajectory features"),
        ("", ""),
        ("+ 4", " embedding features (MiniLM)"),
        ("", " drift, crisis similarity, slope, variance"),
        ("", ""),
        ("+ 3", " coherence features"),
        ("", " short_response, topic_coherence,"),
        ("", " question_response_ratio"),
        ("", ""),
        ("= 67", " total features"),
    ]

    y_calc = Inches(2.05)
    for bold_part, normal_part in breakdown:
        line_box = slide.shapes.add_textbox(Inches(8.5), y_calc, Inches(4.1), Inches(0.3))
        tf = line_box.text_frame
        p = tf.paragraphs[0]
        if bold_part:
            run_b = p.add_run()
            run_b.text = bold_part
            run_b.font.size = Pt(15)
            run_b.font.color.rgb = DARK_BLUE
            run_b.font.bold = True
            run_b.font.name = FONT_NAME
        run_n = p.add_run()
        run_n.text = normal_part
        run_n.font.size = Pt(15)
        run_n.font.color.rgb = TEXT_GRAY
        run_n.font.name = FONT_NAME

        y_calc += Inches(0.28)

    # Highlight the total
    total_box = add_rounded_rect(slide, Inches(8.4), y_calc - Inches(0.38), Inches(4.2), Inches(0.45),
                                  fill_color=ACCENT_BLUE, line_color=ACCENT_BLUE)
    tf = total_box.text_frame
    p = tf.paragraphs[0]
    p.text = "67 features = full emotional trajectory"
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 5)


def build_slide_06_safety_gates(prs):
    """Slide 6: 6-Gate Safety Logic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "6-Gate Safety Logic", "Layered protection: ML cannot override safety rules")

    gates = [
        ("Gate 1", "< 3 user messages", "\u2192 Green (not enough data)", GREEN),
        ("Gate 2", "Safety keyword floor", "\u2192 Force Orange/Red", RED),
        ("Gate 3", "ML prediction", "\u2192 GradientBoosting classification", ACCENT_BLUE),
        ("Gate 4", "Low confidence (< 45%)", "\u2192 Default to Yellow", YELLOW),
        ("Gate 5", "Short conversation (< 6 msgs)", "\u2192 Cap at Orange max", ORANGE),
        ("Gate 6", "Safety floor enforcement", "\u2192 ML can never go below keyword level", RED),
    ]

    y = Inches(1.45)
    for gate_name, condition, action, color in gates:
        # Gate number badge
        badge = add_rounded_rect(slide, Inches(0.7), y, Inches(1.1), Inches(0.6),
                                  fill_color=color, line_color=color)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = gate_name
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Condition
        cond_box = slide.shapes.add_textbox(Inches(2.0), y + Inches(0.05), Inches(4.0), Inches(0.5))
        tf = cond_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = condition
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.font.bold = True
        p.font.name = FONT_NAME

        # Action
        act_box = slide.shapes.add_textbox(Inches(6.0), y + Inches(0.05), Inches(4.5), Inches(0.5))
        tf = act_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = FONT_NAME

        y += Inches(0.78)

    # Quote box
    add_quote_box(slide, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.65),
                  "Asymmetric errors: false positives are always preferable to false negatives.",
                  font_size=16, color=DARK_BLUE)

    # Visual: flow diagram on right side
    flow_box = add_rounded_rect(slide, Inches(10.8), Inches(1.45), Inches(2.0), Inches(4.5),
                                 fill_color=LIGHT_GRAY, line_color=MEDIUM_GRAY)
    flow_title = slide.shapes.add_textbox(Inches(10.9), Inches(1.55), Inches(1.8), Inches(0.35))
    p = flow_title.text_frame.paragraphs[0]
    p.text = "Priority Order"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    priority_items = ["Safety rules\n(always win)", "\u25bc", "ML model\n(when confident)", "\u25bc", "Cautious\ndefaults"]
    y_p = Inches(2.0)
    for item in priority_items:
        tb = slide.shapes.add_textbox(Inches(10.9), y_p, Inches(1.8), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY if "\u25bc" not in item else MEDIUM_GRAY
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        y_p += Inches(0.55) if "\u25bc" not in item else Inches(0.35)

    add_slide_number(slide, 6)


def build_slide_07_adversarial(prs):
    """Slide 7: Track 1 — Adversarial Stress-Testing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Track 1: Adversarial Stress-Testing", "36 test cases across 20 categories")

    # Progress: before/after
    progress_data = [
        ["", "Baseline", "Final"],
        ["Tests passing", "7 / 10", "36 / 36"],
        ["Critical misses", "3", "0"],
        ["Categories", "4", "20"],
    ]
    add_table(slide, Inches(0.7), Inches(1.45), Inches(5.0), progress_data,
              col_widths=[Inches(2.0), Inches(1.5), Inches(1.5)],
              font_size=15, row_height=Inches(0.48))

    # Categories grid (right side)
    cat_title = slide.shapes.add_textbox(Inches(6.3), Inches(1.45), Inches(6.5), Inches(0.35))
    p = cat_title.text_frame.paragraphs[0]
    p.text = "Test Categories"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    categories = [
        "sarcasm", "negation", "code-switching", "Qu\u00e9b\u00e9cois slang",
        "gradual drift", "direct crisis", "hidden intent", "manipulation",
        "somatization", "identity conflict", "sudden escalation", "active bypass",
        "rapid recovery", "cultural FP", "neurodivergent", "emoji-only",
        "repeated word", "short recovery", "long message", "emoji crisis",
    ]

    y_cat = Inches(1.9)
    col = 0
    for i, cat in enumerate(categories):
        x = Inches(6.3) + Inches(col * 3.3)
        badge = add_rounded_rect(slide, x, y_cat, Inches(3.0), Inches(0.32),
                                  fill_color=LIGHT_GRAY, line_color=MEDIUM_GRAY)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = f"  \u2713  {cat}"
        p.font.size = Pt(12)
        p.font.color.rgb = GREEN
        p.font.bold = True
        p.font.name = FONT_NAME

        if col == 0:
            col = 1
        else:
            col = 0
            y_cat += Inches(0.38)

    # Bottom: safety rule
    add_quote_box(slide, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.65),
                  "Exit code 2 = critical miss = blocks merge. Every crisis message must be caught.",
                  font_size=15)

    # Additional context bullets
    items = [
        "Automated CI/CD gate: safety regression blocks deployment",
        "Tolerance-based testing: \u00b12 level tolerance for edge cases",
        "Bilingual coverage: FR + EN test cases in every category",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(3.7), Inches(5.0), Inches(2.2), items, font_size=14)

    add_slide_number(slide, 7)


def build_slide_08_logic_hardening(prs):
    """Slide 8: Track 2 — Logic Hardening."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Track 2: Logic Hardening", "Making detection robust against edge cases")

    hardening_items = [
        ("Negation Handling", '"I can\'t cope" vs "I can cope"\nRegex patterns for FR/EN negation structures',
         RGBColor(0xE8, 0xF5, 0xE9), GREEN),
        ("Identity Conflict Detection", '2SLGBTQ+ distress expressions\n"My family won\'t accept who I am"',
         RGBColor(0xE3, 0xF2, 0xFD), ACCENT_BLUE),
        ("Somatization Flag", 'Physical + emotional co-occurrence\n"My chest hurts" = emotional pain in some cultures',
         RGBColor(0xFD, 0xED, 0xED), RED),
        ("Sentence Embeddings", 'Multilingual MiniLM (paraphrase-multilingual)\nCatches synonyms, paraphrases across FR/EN',
         RGBColor(0xFE, 0xF3, 0xE4), ORANGE),
        ("Behavioral Coherence", 'Withdrawal patterns: short responses, topic shifts\nDisengagement signals across conversation',
         RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x8E, 0x24, 0xAA)),
        ("Word-Boundary Matching", 'Regex \\b prevents false positives\n"morte de rire" no longer matches "mort"',
         LIGHT_GRAY, DARK_GRAY),
    ]

    y = Inches(1.4)
    for i, (title, desc, bg, border) in enumerate(hardening_items):
        col = i % 2
        x = Inches(0.5) + Inches(col * 6.3)

        if col == 0 and i > 0:
            y += Inches(1.65)

        box = add_rounded_rect(slide, x, y, Inches(6.0), Inches(1.5), fill_color=bg, line_color=border)

        t_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.font.bold = True
        p.font.name = FONT_NAME

        d_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.9))
        tf = d_box.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_GRAY
        p2.font.name = FONT_NAME

    add_slide_number(slide, 8)


def build_slide_09_data_augmentation(prs):
    """Slide 9: Track 3 — Synthetic Data Augmentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Track 3: Synthetic Data Augmentation", "600 conversations via Claude Haiku API")

    # Standard data
    std_box = add_rounded_rect(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(2.3),
                                fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)

    std_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.35))
    p = std_title.text_frame.paragraphs[0]
    p.text = "480 Standard Conversations"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    std_items = [
        "60 conversations per class \u00d7 4 classes \u00d7 2 languages",
        "Classes: Green (normal), Yellow (concerning), Orange (distress), Red (crisis)",
        "12 user + 12 assistant messages per conversation",
        "Bilingual: authentic Qu\u00e9b\u00e9cois French + Canadian English",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.95), Inches(5.4), Inches(1.7), std_items, font_size=14)

    # Adversarial data
    adv_box = add_rounded_rect(slide, Inches(6.8), Inches(1.45), Inches(6.0), Inches(2.3),
                                fill_color=RGBColor(0xFD, 0xED, 0xED), line_color=RED)

    adv_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.35))
    p = adv_title.text_frame.paragraphs[0]
    p.text = "120 Adversarial Conversations"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    adv_items = [
        "6 archetypes \u00d7 10 convos \u00d7 2 languages",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.4), adv_items, font_size=14)

    # Archetype list
    archetypes = [
        ("physical_only", "Somatic complaints masking emotional pain"),
        ("sarcasm_distress", "Humor and sarcasm hiding real crisis"),
        ("adversarial_bypass", "Attempts to manipulate/bypass detection"),
        ("identity_distress", "2SLGBTQ+ / cultural identity crisis"),
        ("neurodivergent_flat", "Flat affect mimicking emotional baseline"),
        ("crisis_with_deflection", "Crisis signals with active minimization"),
    ]

    y_arch = Inches(2.45)
    for name, desc in archetypes:
        arch_box = slide.shapes.add_textbox(Inches(7.2), y_arch, Inches(5.3), Inches(0.25))
        tf = arch_box.text_frame
        p = tf.paragraphs[0]
        run_b = p.add_run()
        run_b.text = f"\u2022  {name}: "
        run_b.font.size = Pt(12)
        run_b.font.color.rgb = RED
        run_b.font.bold = True
        run_b.font.name = FONT_NAME
        run_n = p.add_run()
        run_n.text = desc
        run_n.font.size = Pt(12)
        run_n.font.color.rgb = TEXT_GRAY
        run_n.font.name = FONT_NAME
        y_arch += Inches(0.22)

    # Bottom: metrics
    metrics_data = [
        ["Metric", "Before", "After"],
        ["Total conversations", "24", "600"],
        ["Sample:feature ratio", "7.2:1", "9.0:1"],
        ["CV variance", "\u00b126.4%", "\u00b11.6%"],
        ["Quality annotation", "\u2014", "Claude-scored"],
    ]
    add_table(slide, Inches(0.5), Inches(4.2), Inches(5.8), metrics_data,
              col_widths=[Inches(2.6), Inches(1.6), Inches(1.6)],
              font_size=14, row_height=Inches(0.42))

    # Note
    note = slide.shapes.add_textbox(Inches(6.8), Inches(4.3), Inches(5.8), Inches(0.6))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "All data is 100% synthetic \u2014 no real PII (strict hackathon rule).\nBilingual parity: every conversation exists in both FR and EN."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_GRAY
    p.font.italic = True
    p.font.name = FONT_NAME

    add_slide_number(slide, 9)


def build_slide_10_ux_handoff(prs):
    """Slide 10: UX & Warm Handoff."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "UX & Warm Handoff", "From cold referral to accompanied transition")

    # 5-step flow
    steps = [
        ("1", "Empathetic\nAcknowledgment", "Validate feelings,\nno hotline yet", GREEN),
        ("2", "Permission-Based\nTransition", "Ask consent, frame as\n\"upgrade\" not rejection", YELLOW),
        ("3", "Context\nBridge", "Anonymized summary\nfor KHP responder", ORANGE),
        ("4", "Seamless\nConnection", "Text 686868,\nsame modality", RED),
        ("5", "Background\nMonitoring", "CEDD stays active,\nacknowledges return", ACCENT_BLUE),
    ]

    x = Inches(0.4)
    for num, title, desc, color in steps:
        # Step badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), Inches(1.5), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Title
        t_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.05), Inches(1.55), Inches(0.7))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_BLUE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Description
        d_box = slide.shapes.add_textbox(x + Inches(0.05), Inches(2.75), Inches(1.65), Inches(0.75))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if num != "5":
            arr = slide.shapes.add_textbox(x + Inches(1.8), Inches(1.55), Inches(0.5), Inches(0.4))
            p = arr.text_frame.paragraphs[0]
            p.text = "\u25b6"
            p.font.size = Pt(14)
            p.font.color.rgb = MEDIUM_GRAY
            p.alignment = PP_ALIGN.CENTER

        x += Inches(2.4)

    # Simulated counselor box
    alex_box = add_rounded_rect(slide, Inches(0.5), Inches(3.75), Inches(6.0), Inches(1.8),
                                 fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)

    alex_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(5.6), Inches(0.35))
    p = alex_title.text_frame.paragraphs[0]
    p.text = "Simulated Counselor \"Alex\""
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    alex_items = [
        "ASIST-trained KHP counselor persona",
        "Active listening: validates before asking",
        "Short responses (2-4 sentences), one question at a time",
        "Blue gradient bubbles distinguish from AI chatbot",
        "Bilingual (FR with tutoiement, EN)",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(4.25), Inches(5.5), Inches(1.2), alex_items, font_size=13)

    # Research stats (right side)
    research_box = add_rounded_rect(slide, Inches(7.0), Inches(3.75), Inches(5.8), Inches(1.8),
                                     fill_color=LIGHT_GRAY, line_color=MEDIUM_GRAY)

    res_title = slide.shapes.add_textbox(Inches(7.2), Inches(3.85), Inches(5.4), Inches(0.35))
    p = res_title.text_frame.paragraphs[0]
    p.text = "Research Evidence"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    research_items = [
        "71% of youth prefer non-verbal communication (JMIR)",
        "44% of 988 callers abandon before connecting (GSA/OES)",
        "20% seek suicide help via text vs 5% via phone (Ontario)",
        "Compare mode: raw LLM vs CEDD-modulated side-by-side",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(4.25), Inches(5.4), Inches(1.2), research_items, font_size=13)

    add_slide_number(slide, 10)


def build_slide_11_results(prs):
    """Slide 11: Results & Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Results & Metrics", "Improvement trajectory: baseline to final")

    # Main comparison table
    results_data = [
        ["Metric", "Baseline", "Final"],
        ["CV Accuracy", "66.7% \u00b1 26.4%", "90.0% \u00b1 1.6%"],
        ["Features", "42", "67"],
        ["Training data", "24 convos", "600 convos"],
        ["Adversarial tests", "7 / 10", "36 / 36"],
        ["Critical misses", "3", "0"],
        ["Languages", "FR only", "FR + EN"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(6.0), results_data,
              col_widths=[Inches(2.2), Inches(1.9), Inches(1.9)],
              font_size=15, row_height=Inches(0.48))

    # Top features box
    feat_box = add_rounded_rect(slide, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.75),
                                 fill_color=LIGHT_GRAY, line_color=ACCENT_BLUE)

    feat_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(5.6), Inches(0.35))
    p = feat_title.text_frame.paragraphs[0]
    p.text = "Top Features by Importance"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    top_feats = [
        ("word_count_max", "0.192"),
        ("word_count_slope", "0.179"),
        ("word_count_last", "0.138"),
        ("length_delta_mean", "0.075"),
    ]

    y_f = Inches(5.5)
    for feat, score in top_feats:
        # Feature name
        f_box = slide.shapes.add_textbox(Inches(0.8), y_f, Inches(3.0), Inches(0.3))
        p = f_box.text_frame.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = FONT_NAME
        p.font.bold = True

        # Score bar (visual)
        bar_width = float(score) * Inches(8)
        add_shape(slide, Inches(3.8), y_f + Inches(0.05), bar_width, Inches(0.2),
                  fill_color=ACCENT_BLUE)

        # Score label
        s_box = slide.shapes.add_textbox(Inches(3.8) + bar_width + Inches(0.1), y_f, Inches(0.6), Inches(0.3))
        p = s_box.text_frame.paragraphs[0]
        p.text = score
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_BLUE
        p.font.bold = True
        p.font.name = FONT_NAME

        y_f += Inches(0.3)

    # vs EmoAgent comparison (right side)
    emo_box = add_rounded_rect(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.3),
                                fill_color=LIGHT_GRAY, line_color=DARK_BLUE)

    emo_title = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.4))
    p = emo_title.text_frame.paragraphs[0]
    p.text = "CEDD vs EmoAgent (Princeton/Michigan)"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    emo_data = [
        ["", "CEDD", "EmoAgent"],
        ["Detection", "~0ms, $0\n(ML only)", "4\u00d7 GPT-4o calls\n(slow, expensive)"],
        ["Explainability", "Full\n(67 named features)", "Black box"],
        ["Languages", "FR + EN native", "English only"],
        ["Cross-session", "SQLite tracking", "Per-convo only"],
        ["Cultural", "Somatization,\nidentity conflict", "None"],
    ]
    add_table(slide, Inches(7.2), Inches(2.0), Inches(5.4), emo_data,
              col_widths=[Inches(1.5), Inches(1.95), Inches(1.95)],
              font_size=12, row_height=Inches(0.62))

    add_slide_number(slide, 11)


def build_slide_12_impact(prs):
    """Slide 12: Impact & Next Steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Impact & Next Steps")

    # Target applications (left)
    target_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35))
    p = target_title.text_frame.paragraphs[0]
    p.text = "Target Applications"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    target_items = [
        "AI emotional support platforms",
        "School chatbots and wellness apps",
        "Kids Help Phone integration",
        "Any youth-facing conversational AI",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.8), Inches(5.5), Inches(1.5), target_items, font_size=15)

    # Canadian multicultural advantage
    multi_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(5.5), Inches(0.35))
    p = multi_title.text_frame.paragraphs[0]
    p.text = "Canadian Multicultural Advantage"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    multi_items = [
        "Somatization detection (South Asian cultural patterns)",
        "Identity conflict detection (2SLGBTQ+ youth)",
        "Withdrawal/silence detection (East Asian expression norms)",
        "Native French + English bilingual support",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(3.7), Inches(5.5), Inches(1.5), multi_items, font_size=14)

    # Next steps (right)
    next_box = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.5),
                                 fill_color=LIGHT_GRAY, line_color=ACCENT_BLUE)

    next_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.35))
    p = next_title.text_frame.paragraphs[0]
    p.text = "Next Steps"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_NAME

    next_items = [
        "LSTM sequence model (understand message order natively)",
        "Real clinical data partnership for validation",
        "PHQ-9 integration as complementary metric",
        "REST API for plug-in deployment",
        "Indigenous expression pattern expansion",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(1.8), next_items, font_size=14)

    # Emergency resources box
    resource_box = add_rounded_rect(slide, Inches(6.8), Inches(4.2), Inches(6.0), Inches(1.8),
                                     fill_color=RGBColor(0xFD, 0xED, 0xED), line_color=RED)

    res_title = slide.shapes.add_textbox(Inches(7.0), Inches(4.3), Inches(5.6), Inches(0.35))
    p = res_title.text_frame.paragraphs[0]
    p.text = "Emergency Resources"
    p.font.size = Pt(16)
    p.font.color.rgb = RED
    p.font.bold = True
    p.font.name = FONT_NAME

    resources = [
        "Kids Help Phone: 1-800-668-6868",
        "Text: 686868",
        "Suicide Crisis Helpline: 9-8-8",
        "Emergency: 911",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(4.65), Inches(5.6), Inches(1.2), resources, font_size=14, color=RED)

    # Closing quote
    add_quote_box(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.8),
                  "CEDD detects the drift, not just the alarm cry.",
                  font_size=22, color=DARK_BLUE)

    # Team footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Team 404HarmNotFound  \u00b7  Shuchita Singh  \u00b7  Amanda Wu  \u00b7  Priyanka Naga  \u00b7  Dominic D'Apice"
    p.font.size = Pt(14)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 12)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    print("Generating presentation...")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_01_title(prs)
    print("  [1/12] Title slide")

    build_slide_02_problem(prs)
    print("  [2/12] The Problem")

    build_slide_03_solution(prs)
    print("  [3/12] Our Solution")

    build_slide_04_architecture(prs)
    print("  [4/12] Architecture")

    build_slide_05_features(prs)
    print("  [5/12] Feature Engineering")

    build_slide_06_safety_gates(prs)
    print("  [6/12] 6-Gate Safety Logic")

    build_slide_07_adversarial(prs)
    print("  [7/12] Adversarial Stress-Testing")

    build_slide_08_logic_hardening(prs)
    print("  [8/12] Logic Hardening")

    build_slide_09_data_augmentation(prs)
    print("  [9/12] Synthetic Data Augmentation")

    build_slide_10_ux_handoff(prs)
    print("  [10/12] UX & Warm Handoff")

    build_slide_11_results(prs)
    print("  [11/12] Results & Metrics")

    build_slide_12_impact(prs)
    print("  [12/12] Impact & Next Steps")

    prs.save(OUTPUT_FILE)
    print(f"\nPresentation saved to: {OUTPUT_FILE}")
    print("Done!")


if __name__ == "__main__":
    main()
